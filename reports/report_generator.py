import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import sys

# Usage: python reports/report_generator.py data/filtered_air_quality.csv

CSV_FILE = sys.argv[1] if len(sys.argv) > 1 else "data/filtered_air_quality.csv"
IMG_FILE = "reports/chart.png"
PPTX_FILE = "reports/air_quality_report.pptx"

# Load data
df = pd.read_csv(CSV_FILE, parse_dates=["datetime"])

# Create a chart (NO2 over time)
plt.figure(figsize=(10,4))
plt.plot(df["datetime"], df["no2"], label="NO₂")
plt.xlabel("Datetime")
plt.ylabel("NO₂ (µg/m³)")
plt.title("NO₂ Time Series")
plt.tight_layout()
plt.savefig(IMG_FILE)
plt.close()

# Create PPTX
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.add_picture(IMG_FILE, Inches(1), Inches(1), width=Inches(8), height=Inches(4))
prs.save(PPTX_FILE)

print(f"Report created: {PPTX_FILE}")
