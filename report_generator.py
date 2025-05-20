# report_generator.py

import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF
import os

# === USER INPUTS ===
CSV_FILE = "filtered_air_quality.csv"  # Export this from your Streamlit app
CHART_PNG = "timeseries.png"

# --- 1. Load CSV and Plot Chart ---
df = pd.read_csv(CSV_FILE, parse_dates=["datetime"])

# Choose which metrics to plot
metrics = [col for col in ["no2", "pm25"] if col in df.columns]
if not metrics:
    raise ValueError("CSV does not contain NO₂ or PM₂.₅ columns!")

fig = px.line(
    df,
    x="datetime",
    y=metrics,
    color="site_name" if "site_name" in df.columns else None,
    title="Time Series: " + " & ".join([m.upper() for m in metrics]),
    labels={"value": "Measurement", "variable": "Metric"}
)
fig.write_image(CHART_PNG)
print(f"✅ Chart image saved as {CHART_PNG}")

# --- 2. Create PowerPoint ---
pptx_file = "air_quality_report.pptx"
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.add_picture(CHART_PNG, Inches(1), Inches(1), width=Inches(8))
prs.save(pptx_file)
print(f"✅ PowerPoint saved as {pptx_file}")

# --- 3. Create PDF ---
pdf_file = "air_quality_report.pdf"
pdf = FPDF()
pdf.add_page()
pdf.image(CHART_PNG, x=10, y=20, w=180)
pdf.set_font("Arial", size=12)
pdf.ln(100)
if os.path.exists(CSV_FILE):
    pdf.cell(0, 10, f"See attached data file: {CSV_FILE}", ln=True)
pdf.output(pdf_file)
print(f"✅ PDF saved as {pdf_file}")

print("Done! Attach the CSV if needed when sharing the report.")

# --- Dependencies required:
# pip install pandas plotly kaleido python-pptx fpdf
